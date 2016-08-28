# -*- coding: utf-8 -*-
from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches, Pt
import argparse
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
import pandas as pd
import numpy as np
from pptx.dml.color import RGBColor
import calendar
import datetime

# P1：总体环比情况表
def createMomTable(mom_in_df, mom_out_df):
    rows = 3
    cols = 6
    left = Inches(0.2)
    top = Inches(1.5)
    width = Inches(9.6)
    height = Inches(2.4)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(1.6)
    table.columns[4].width = Inches(1.6)
    table.columns[5].width = Inches(1.6)

    # write column headings
    table.cell(0, 0).text = '漫游方向'
    table.cell(0, 1).text = '总话单量（亿条）'
    table.cell(0, 2).text = '总用户数（百万户）'
    table.cell(0, 3).text = '通话时长（千万分钟）'
    table.cell(0, 4).text = '短信条数（亿条）'
    table.cell(0, 5).text = '流量（TB）'

    # write body cells
    table.cell(1, 0).text = '来访'
    table.cell(1, 1).text = str(mom_in_df.COUNTCDR.values[0]) + '\r'+ str(mom_in_df.COUNTCDRPCT.values[0]) + '%'
    table.cell(1, 2).text = str(mom_in_df.USERCNT.values[0]) + '\r'+ str(mom_in_df.USERPCT.values[0]) + '%'
    table.cell(1, 3).text = str(mom_in_df.CALLDURATION.values[0]) + '\r'+ str(mom_in_df.CALLDURATIONPCT.values[0]) + '%'
    table.cell(1, 4).text = str(mom_in_df.SMS.values[0]) + '\r'+ str(mom_in_df.SMSPCT.values[0]) + '%'
    table.cell(1, 5).text = str(mom_in_df.DATAALL.values[0]) + '\r'+ str(mom_in_df.DATAALLPCT.values[0]) + '%'

    # write body cells
    table.cell(2, 0).text = '出访'
    table.cell(2, 1).text = str(mom_out_df.COUNTCDR.values[0]) + '\r'+ str(mom_out_df.COUNTCDRPCT.values[0]) + '%'
    table.cell(2, 2).text = str(mom_out_df.USERCNT.values[0]) + '\r'+ str(mom_out_df.USERPCT.values[0]) + '%'
    table.cell(2, 3).text = str(mom_out_df.CALLDURATION.values[0]) + '\r'+ str(mom_out_df.CALLDURATIONPCT.values[0]) + '%'
    table.cell(2, 4).text = str(mom_out_df.SMS.values[0]) + '\r'+ str(mom_out_df.SMSPCT.values[0]) + '%'
    table.cell(2, 5).text = str(mom_out_df.DATAALL.values[0]) + '\r'+ str(mom_out_df.DATAALLPCT.values[0]) + '%'

# create cluster-bar-char
def createTrendChart(chart_data, x, y, cx, cy):

    # add chart to slide1--------------------

    graphic_frame = slide.shapes.add_chart(
        # XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart

    #legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    # category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(8)

    value_axis = chart.value_axis
    # value_axis.maximum_scale = 50.0
    # value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.has_major_gridlines = False
    value_axis.tick_labels.font.size = Pt(8)

# create line-char
def createLineChart(slide, chart_data, x, y, cx, cy):

    # add chart to slide
    graphic_frame = slide.shapes.add_chart(
        # XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart

    #legend
    chart.has_legend = False
    # chart.legend.position = XL_LEGEND_POSITION.RIGHT
    # chart.legend.include_in_layout = False
    # chart.legend.font.size = Pt(9)

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    # category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(8)

    value_axis = chart.value_axis
    # value_axis.maximum_scale = 50.0
    # value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.has_major_gridlines = False
    value_axis.tick_labels.font.size = Pt(8)

# create pie-char
def createPieChart(slide, chart_data, x, y, cx, cy):

    # add chart to slide
    graphic_frame = slide.shapes.add_chart(
        # XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart

    #legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.font.size = Pt(9)

# create slide 1.x
def createmomana(pagetitle,titlenm, unit, MONTH, ROAM_TYPE, BUSI_TYPE, bv):
    # slide1.1 来访话单量环比情况
    title_only_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    # slide1.1 chart1 per day
    title = ''
    pagetitlepls = slide.placeholders[14]
    pagetitlepls.text = pagetitle
    titlepls = slide.placeholders[13]
    titlepls.text = titlenm
    unitpls = slide.placeholders[12]
    unitpls.text = unit

    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 16.7, 19.2, 21.4, 16.7, 12.0)
    doddf = pd.read_excel("E:\djangoWeb\pptgen\dod_pptgen.xlsx").sort_values(by='CALLDATE')

    dodin = doddf[
        ['COUNTCDR','USERCNT','CALLDURATION','SMS','DATAALL']][( doddf['MONTH'] == int(MONTH_ANA) ) & (doddf.ROAM_TYPE == ROAM_TYPE)]
    valuetuple = dodin[bv].values
    # print(dodin)
    # define chart data
    chart_data = ChartData()
    chart_data.categories = SILDE1x_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.7), Inches(10), Inches(2.8)
    createLineChart(slide, chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per carrier
    # piechar
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distcarriercdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_carrier.xlsx")
    distcarrierin = distcarriercdrdf[
        ['CARRIER_CD', 'PERCENT']][
        (distcarriercdrdf.MONTH == int(MONTH)) & (distcarriercdrdf.ROAM_TYPE == ROAM_TYPE) & (distcarriercdrdf.BUSI_TYPE == BUSI_TYPE)]
    valuetuple = distcarrierin.PERCENT.values
    valuetuplefirst5 = distcarrierin[distcarrierin.CARRIER_CD != 'Others'].PERCENT.values

    # first5 PERCENT
    incdrfirstcarrierpls = slide.placeholders[10]
    incdrfirstcarrierpls.text = str(valuetuplefirst5.sum()*100)+'%'
    # incdrfirstcarrierpls.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # valuetuple.resize(6)
    # otherpct = 1 - valuetuple.sum()
    # valuetuple[5]=otherpct
    # valuetuple = valuetuple.tolist().append(otherpct)
    # print(valuetuple)
    # print(list(distcarrierin.CARRIER_CD.values).append('Others'))
    # define chart data ---------------------
    chart_data = ChartData()
    # carrierCategories = ('HKGPP', 'TWNFE', 'MACCT', 'USACG', 'THAWP', 'Others')
    carrierCategories = distcarrierin.CARRIER_CD.values
    chart_data.categories = carrierCategories
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(slide, chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per prov
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distprovcdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_prov.xlsx")
    distprovin = distprovcdrdf[
        ['ENGNM', 'PERCENT']][
        (distprovcdrdf.MONTH == int(MONTH)) & (distprovcdrdf.ROAM_TYPE == ROAM_TYPE) & (
            distprovcdrdf.BUSI_TYPE == BUSI_TYPE)]
    valuetuple = distprovin.PERCENT.values
    valuetuplefirst5 = distprovin[distprovin.ENGNM != 'Others'].PERCENT.values

    # first5 PERCENT
    incdrfirstprovpls = slide.placeholders[11]
    incdrfirstprovpls.text = str(valuetuplefirst5.sum()*100)+'%'

    # define chart data
    chart_data = ChartData()
    provCategories = ('BJ', 'GD', 'SH', 'JS', 'ZJ', 'Others')
    # provCategories = distprovin.ENGNM.values
    chart_data.categories = provCategories

    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(slide, chart_data, x, y, cx, cy)

def getmonthlist(v_month):
    result = []
    month_tmp = v_month
    one_day = datetime.timedelta(days=1)
    for x in range(1,14):
        result.append(month_tmp)
        lastmonth_lastdate = datetime.datetime.strptime(month_tmp, '%Y%m') - one_day
        # get last month str
        month_tmp = datetime.datetime.strftime(lastmonth_lastdate, '%Y%m')
    return result

# 函数在这里运行

if __name__ == "__main__":

    # constant variable
    MONTH_ANA = '201605'
    # Get month list for categories
    MONTH_LIST = getmonthlist(MONTH_ANA)
    MONTH_LIST.reverse()
    # Get days
    DAYS = calendar.monthrange(int(MONTH_ANA[0:4]), int(MONTH_ANA[4:6]))[1]

    # SILDE1_CATEGORIES = ['201506', '201508', '201509', '201510', '201511', '201512', '201601', '201602', '201603', '201604', '201605', '201606', '201606']
    SILDE1_CATEGORIES = MONTH_LIST
    # SILDE1x_CATEGORIES = ['1', '2', '3', '4', '5', '6', '7', '8', '9', '10', '11', '12', '13', '14', '15', '16', '17', '18', '19', '20', '21', '22', '23', '24', '25', '26', '27', '28', '29', '30']
    SILDE1x_CATEGORIES = [x for x in range(1, DAYS+1)]

    prs = Presentation('ana-template.pptx')

    # slide1 总体环比情况
    title_only_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    # table = slide.placeholders[11]
    paragraph1 = slide.placeholders[12]
    paragraph2 = slide.placeholders[13]
    paragraph1.text = "来访业务量总体有所下降，其中，通话时长降幅较大"
    paragraph2.text = "出访业务量总体有所下降，其中，通话时长降幅较大"
    # slide1 table 环比表格
    # 获取数据
    df = pd.read_excel("E:\djangoWeb\pptgen\mom_pptgen.xlsx").sort_values(by='MONTH')
    try:
        mom_in_df = df[
            ['COUNTCDR', 'USERCNT', 'CALLDURATION', 'SMS', 'DATAALL', 'COUNTCDRPCT', 'USERPCT', 'CALLDURATIONPCT', 'SMSPCT',
             'DATAALLPCT', 'USER4GPCT', 'DOU4GPCT']][(df.MONTH == int(MONTH_ANA)) & (df.ROAM_TYPE == 1)]
        mom_out_df = df[
            ['COUNTCDR', 'USERCNT', 'CALLDURATION', 'SMS', 'DATAALL', 'COUNTCDRPCT', 'USERPCT', 'CALLDURATIONPCT', 'SMSPCT',
             'DATAALLPCT', 'USER4GPCT', 'DOU4GPCT']][(df.MONTH == int(MONTH_ANA)) & (df.ROAM_TYPE == 2)]
    except Exception as e:
        raise e
    # print(mom_in_df)
    createMomTable(mom_in_df, mom_out_df)

    # slide1.1 来访话单量环比情况
    # createmomana(pagetitle, titlenm, unit, MONTH, ROAM_TYPE, BUSI_TYPE, bv)
    pagetitle = u"工作总结——生产数据情况-环比-来访话单数"
    titlenm = u"来访话单量按日情况"
    unit = u"单位：百万条"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 1, 1, 'COUNTCDR')

    # slide1.2 来访用户数环比情况
    pagetitle = u"工作总结——生产数据情况-环比-用户数"
    titlenm = u"来访用户数按日情况"
    unit = u"单位：万户"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 1, 2, 'USERCNT')

    # slide1.3 来访通话时长环比情况
    pagetitle = u"工作总结——生产数据情况-环比-来访通话时长"
    titlenm = u"来访通话时长按日情况"
    unit = u"单位：万分钟"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 1, 3, 'CALLDURATION')

    # slide1.4 来访短信环比情况
    pagetitle = u"工作总结——生产数据情况-环比-来访短信"
    titlenm = u"来访短信按日情况"
    unit = u"单位：万条"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 1, 4, 'SMS')

    # slide1.5 来访数据环比情况
    pagetitle = u"工作总结——生产数据情况-环比-来访数据"
    titlenm = u"来访数据按日情况"
    unit = u"单位：TB"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 1, 5, 'DATAALL')

    # slide1.6 出访话单量环比情况
    pagetitle = u"工作总结——生产数据情况-环比-出访话单数"
    titlenm = u"出访话单量按日情况"
    unit = u"单位：百万条"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 2, 1, 'COUNTCDR')

    # slide1.7 出访用户数环比情况
    pagetitle = u"工作总结——生产数据情况-环比-出访用户数"
    titlenm = u"出访用户数按日情况"
    unit = u"单位：万户"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 2, 2, 'USERCNT')

    # slide1.8 出访通话时长环比情况
    pagetitle = u"工作总结——生产数据情况-环比-出访通话时长"
    titlenm = u"出访通话时长按日情况"
    unit = u"单位：万分钟"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 2, 3, 'CALLDURATION')

    # slide1.9 出访短信环比情况
    pagetitle = u"工作总结——生产数据情况-环比-出访短信"
    titlenm = u"出访短信按日情况"
    unit = u"单位：万条"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 2, 4, 'SMS')

    # slide1.10 出访数据环比情况
    pagetitle = u"工作总结——生产数据情况-环比-出访数据"
    titlenm = u"出访数据按日情况"
    unit = u"单位：TB"
    createmomana(pagetitle, titlenm, unit, MONTH_ANA, 2, 5, 'DATAALL')

    # slide2 传统业务量趋势情况 - 来访
    title_only_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    momdf = pd.read_excel("E:\djangoWeb\pptgen\mom_pptgen.xlsx").sort_values(by='MONTH')
    momin = momdf[
        ['COUNTCDR', 'USERCNT', 'CALLDURATION', 'SMS', 'DATAALL', 'DATA23G', 'DATA4G', 'USER4G',  'DOU4G']][
        (momdf.MONTH <= int(MONTH_ANA)) & (momdf.MONTH >= (int(MONTH_ANA)-100)) & (momdf.ROAM_TYPE == 1)]

    # slide2 chart1 CDR
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple = momin['COUNTCDR'].values

    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.5), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide2 chart2 CALLDURATION
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple = momin['CALLDURATION'].values
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(2.7), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide2 chart3 SMS
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple = momin['SMS'].values
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(4.9), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide3 传统业务量趋势情况 - 出访
    title_only_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    momout = momdf[
        ['COUNTCDR', 'USERCNT', 'CALLDURATION', 'SMS', 'DATAALL', 'DATA23G', 'DATA4G', 'USER4G', 'DOU4G']][
        (momdf.MONTH <= int(MONTH_ANA)) & (momdf.MONTH >= (int(MONTH_ANA) - 100)) & (
        momdf.ROAM_TYPE == 2)]

    # slide3 chart1 CDR
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple = momout['COUNTCDR'].values
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.5), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide3 chart2 CALLDURATION
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple = momout['CALLDURATION'].values
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(2.7), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide3 chart3 SMS
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple = momout['SMS'].values
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(4.9), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 数据业务量趋势情况
    title_only_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    # momindata = momdf[
    #     ['DATAALL', 'data23g', 'data4g']][
    #     (momdf.MONTH.values[0] <= int(MONTH_ANA)) & (momdf.MONTH.values[0] >= (int(MONTH_ANA)-100)) & (momdf.ROAM_TYPE == 1)]

    # slide4 chart1 来访
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple23g = momin['DATA23G'].values
    valuetuple4g = momin['DATA4G'].values
    # define chart data1 ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("23g", valuetuple23g)
    chart_data.add_series("4g", valuetuple4g)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0.1), Inches(1.2), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 chart2 出访
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetuple23g = momout['DATA23G'].values
    valuetuple4g = momout['DATA4G'].values
    # define chart data1 ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("23g", valuetuple23g)
    chart_data.add_series("4g", valuetuple4g)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(1.2), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 chart3 用户数
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetupleuser4gin = momin['USER4G'].values
    valuetupleuser4gout = momout['USER4G'].values
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("In", valuetupleuser4gin)
    chart_data.add_series("Out", valuetupleuser4gout)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0.1), Inches(4.7), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 chart4 DOU
    title = ""
    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    valuetupledou4gin = momin['DOU4G'].values
    valuetupledou4gout = momout['DOU4G'].values
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("In", valuetupledou4gin)
    chart_data.add_series("Out", valuetupledou4gout)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(4.7), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 paragraph
    # inroam4gcarriersnuminc = 4
    # inroam4gcarriersnumall = 103
    # outroam4gcarriersnuminc = 4
    # outroam4gcarriersnumall = 122

    # 获取值
    carrieropendf = pd.read_excel("E:\djangoWeb\pptgen\carrier_open_4g.xlsx")
    carrieropenin = carrieropendf[
        ['INCR_NM', 'ACC_NM']][
        (carrieropendf.MONTH == int(MONTH_ANA)) & (carrieropendf.ROAM_TYPE == 1)]
    carrieropenout = carrieropendf[
        ['INCR_NM', 'ACC_NM']][
        (carrieropendf.MONTH == int(MONTH_ANA)) & (carrieropendf.ROAM_TYPE == 2)]

    inroam4gcarriersnuminc = carrieropenin.INCR_NM.values[0]
    inroam4gcarriersnumall = carrieropenin.ACC_NM.values[0]
    outroam4gcarriersnuminc = carrieropenout.INCR_NM.values[0]
    outroam4gcarriersnumall = carrieropenout.ACC_NM.values[0]

    para1 = slide.placeholders[15]
    para1.text = "2016年02月来访LTE新开通" + str(inroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(inroam4gcarriersnumall) + "家；" + "\r" + "2016年02月出访LTE新开通" + str(outroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(outroam4gcarriersnumall) + "家；"

    valuetupleuser4ginpct = mom_in_df['USER4GPCT'].values[0]
    valuetupledou4ginpct = mom_in_df['DOU4GPCT'].values[0]
    valuetupleuser4goutpct = mom_out_df['USER4GPCT'].values[0]
    valuetupledou4goutpct = mom_out_df['DOU4GPCT'].values[0]
    para2 = slide.placeholders[16]
    para2.text = "来访："+ str(valuetupleuser4ginpct) +"%   出访：" + str(valuetupleuser4goutpct) + "%"
    para3 = slide.placeholders[17]
    para3.text = "来访："+ str(valuetupledou4ginpct) +"%   出访：" + str(valuetupledou4goutpct) + "%"

    # tf = para.text_frame
    # p = tf.add_paragraph()
    # p.text = "2016年02月来访LTE新开通" + str(inroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(inroam4gcarriersnumall) + "家；"
    # p.level = 1
    #
    # p = tf.add_paragraph()
    # p.text = "2016年02月出访LTE新开通" + str(outroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(outroam4gcarriersnumall) + "家；"
    # p.level = 1

    prs.save('report1.pptx')
