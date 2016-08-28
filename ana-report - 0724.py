# -*- coding: utf-8 -*-
from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches, Pt
import argparse
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
import pandas as pd
import numpy as np
from pptx.dml.color import RGBColor

# P1：总体环比情况表
def createMomTable(month):
    rows = 3
    cols = 6
    left = Inches(0.2)
    top = Inches(1.5)
    width = Inches(9.6)
    height = Inches(2.4)

    table = shapes.add_table(rows, cols, left, top, width, height).table
    df = pd.read_excel("E:\djangoWeb\pptgen\mom_pptgen.xlsx")
    try:
        mom_in_df=df[['count','user','callduration','sms','dataall','countpct','userpct','calldurationpct','smspct','dataallpct']][(df.month==month)&(df.roam_type==1)]
        mom_out_df=df[['count','user','callduration','sms','dataall','countpct','userpct','calldurationpct','smspct','dataallpct']][(df.month==month)&(df.roam_type==2)]
    except Exception as e:
        raise e
    # set column widths
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(1.6)
    table.columns[4].width = Inches(1.6)
    table.columns[5].width = Inches(1.6)

    # write column headings
    table.cell(0, 0).text = '漫游方向'
    table.cell(0, 1).text = '总话单量（亿条）'
    table.cell(0, 2).text = '总用户数（百万户）'
    table.cell(0, 3).text = '通话时长（千万分钟）'
    table.cell(0, 4).text = '短信条数（亿条）'
    table.cell(0, 5).text = '流量（TB）'

    # write body cells
    table.cell(1, 0).text = '来访'
    table.cell(1, 1).text = str(mom_in_df['count'].values[0]) + '  '+ str(mom_in_df['countpct'].values[0]) + '%'
    table.cell(1, 2).text = str(mom_in_df.user.values[0]) + '  '+ str(mom_in_df.userpct.values[0]) + '%'
    table.cell(1, 3).text = str(mom_in_df.callduration.values[0]) + '  '+ str(mom_in_df.calldurationpct.values[0]) + '%'
    table.cell(1, 4).text = str(mom_in_df.sms.values[0]) + '  '+ str(mom_in_df.smspct.values[0]) + '%'
    table.cell(1, 5).text = str(mom_in_df.dataall.values[0]) + '  '+ str(mom_in_df.dataallpct.values[0]) + '%'

    # write body cells
    table.cell(2, 0).text = '出访'
    table.cell(2, 1).text = str(mom_out_df['count'].values[0]) + '  '+ str(mom_out_df['countpct'].values[0]) + '%'
    table.cell(2, 2).text = str(mom_out_df.user.values[0]) + '  '+ str(mom_out_df.userpct.values[0]) + '%'
    table.cell(2, 3).text = str(mom_out_df.callduration.values[0]) + '  '+ str(mom_out_df.calldurationpct.values[0]) + '%'
    table.cell(2, 4).text = str(mom_out_df.sms.values[0]) + '  '+ str(mom_out_df.smspct.values[0]) + '%'
    table.cell(2, 5).text = str(mom_out_df.dataall.values[0]) + '  '+ str(mom_out_df.dataallpct.values[0]) + '%'

# create cluster-bar-char
def createTrendChart(chart_data, x, y, cx, cy):

    # add chart to slide1--------------------

    graphic_frame = slide.shapes.add_chart(
        # XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart

    #legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    # category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(8)

    value_axis = chart.value_axis
    # value_axis.maximum_scale = 50.0
    # value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.has_major_gridlines = False
    value_axis.tick_labels.font.size = Pt(8)

# create line-char
def createLineChart(chart_data, x, y, cx, cy):

    # add chart to slide
    graphic_frame = slide.shapes.add_chart(
        # XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart

    #legend
    chart.has_legend = False
    # chart.legend.position = XL_LEGEND_POSITION.RIGHT
    # chart.legend.include_in_layout = False
    # chart.legend.font.size = Pt(9)

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    # category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(8)

    value_axis = chart.value_axis
    # value_axis.maximum_scale = 50.0
    # value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.has_major_gridlines = False
    value_axis.tick_labels.font.size = Pt(8)

# create pie-char
def createPieChart(chart_data, x, y, cx, cy):

    # add chart to slide
    graphic_frame = slide.shapes.add_chart(
        # XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart

    #legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.font.size = Pt(9)

# create slide 1.x
def createmomana(titlenm, unit, month, roam_type, busi_type, bv):
    # slide1.1 来访话单量环比情况
    title_only_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    # slide1.1 chart1 per day
    title = ''
    titlenm = u"来访话单量按日情况"
    unit = u"单位：百万条"
    titlepls = slide.placeholders[13]
    titlepls.text = titlenm
    unitpls = slide.placeholders[12]
    unitpls.text = unit

    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 16.7, 19.2, 21.4, 16.7, 12.0)
    doddf = pd.read_excel("E:\djangoWeb\pptgen\dod_pptgen.xlsx")
    dodin = doddf[
        ['count','user','callduration','sms','dataall']][(str(doddf.date.values[0])[0:6] == month) & (doddf.roam_type == roam_type)]
    valuetuple = dodin[bv].values
    # define chart data
    chart_data = ChartData()
    chart_data.categories = SILDE1x_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.7), Inches(10), Inches(2.8)
    createLineChart(chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per carrier
    # piechar
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distcarriercdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_carrier.xlsx")
    distcarrierin = distcarriercdrdf[
        ['carrier_cd', 'percent']][
        (str(distcarriercdrdf.month.values[0]) == month) & (distcarriercdrdf.roam_type == roam_type) & (distcarriercdrdf.busi_type == busi_type)]
    valuetuple = distcarrierin.percent.values
    valuetuplefirst5 = distcarrierin[distcarrierin.carrier_cd != 'Others'].percent.values

    # first5 percent
    incdrfirstcarrierpls = slide.placeholders[10]
    incdrfirstcarrierpls.text = str(valuetuplefirst5.sum()*100)+'%'
    # incdrfirstcarrierpls.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # valuetuple.resize(6)
    # otherpct = 1 - valuetuple.sum()
    # valuetuple[5]=otherpct
    # valuetuple = valuetuple.tolist().append(otherpct)
    # print(valuetuple)
    # print(list(distcarrierin.carrier_cd.values).append('Others'))
    # define chart data ---------------------
    chart_data = ChartData()
    # carrierCategories = ('HKGPP', 'TWNFE', 'MACCT', 'USACG', 'THAWP', 'Others')
    carrierCategories = distcarrierin.carrier_cd.values
    chart_data.categories = carrierCategories
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per prov
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distprovcdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_prov.xlsx")
    distprovin = distprovcdrdf[
        ['prov_cd', 'percent']][
        (str(distprovcdrdf.month.values[0]) == month) & (distprovcdrdf.roam_type == roam_type) & (
            distprovcdrdf.busi_type == busi_type)]
    valuetuple = distprovin.percent.values
    valuetuplefirst5 = distprovin[distprovin.prov_cd != 'Others'].percent.values

    # first5 percent
    incdrfirstprovpls = slide.placeholders[11]
    incdrfirstprovpls.text = str(valuetuplefirst5.sum()*100)+'%'

    # define chart data
    chart_data = ChartData()
    provCategories = ('BJ', 'GD', 'SH', 'JS', 'ZJ', 'Others')
    # provCategories = distprovin.prov_cd.values
    chart_data.categories = provCategories

    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(chart_data, x, y, cx, cy)


# 函数在这里运行

if __name__ == "__main__":

    # constant variable
    # slide1 categories
    MONTH_ANA = '201607'
    SILDE1_CATEGORIES = ['201507', '201508', '201509', '201510', '201511', '201512', '201601', '201602', '201603', '201604', '201605', '201606', '201607']
    SILDE1x_CATEGORIES = ['1', '2', '3', '4', '5', '6', '7', '8', '9', '10', '11', '12', '13', '14', '15', '16', '17', '18', '19', '20', '21', '22', '23', '24', '25', '26', '27', '28', '29', '30', '31']

    prs = Presentation('ana-template.pptx')

    # slide1 总体环比情况
    title_only_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    # table = slide.placeholders[11]
    paragraph1 = slide.placeholders[12]
    paragraph2 = slide.placeholders[13]
    paragraph1.text = "来访业务量总体有所下降，其中，通话时长降幅较大"
    paragraph2.text = "出访业务量总体有所下降，其中，通话时长降幅较大"
    # slide1 table 环比表格
    createMomTable(201607)

    # slide1.1 来访话单量环比情况
    # createmomana(titlenm, unit, month, roam_type, busi_type, bv)

    title_only_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    # slide1.1 chart1 per day
    title = ''
    titlenm = u"来访话单量按日情况"
    unit = u"单位：百万条"
    titlepls = slide.placeholders[13]
    titlepls.text = titlenm
    unitpls = slide.placeholders[12]
    unitpls.text = unit

    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 16.7, 19.2, 21.4, 16.7, 12.0)
    doddf = pd.read_excel("E:\djangoWeb\pptgen\dod_pptgen.xlsx")
    dodin = doddf[
        ['count','user','callduration','sms','dataall']][(str(doddf.date.values[0])[0:6] == MONTH_ANA) & (doddf.roam_type == 1)]
    valuetuple = dodin['count'].values
    # define chart data
    chart_data = ChartData()
    chart_data.categories = SILDE1x_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.7), Inches(10), Inches(2.8)
    createLineChart(chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per carrier
    # piechar
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distcarriercdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_carrier.xlsx")
    distcarrierin = distcarriercdrdf[
        ['carrier_cd', 'percent']][
        (str(distcarriercdrdf.month.values[0]) == MONTH_ANA) & (distcarriercdrdf.roam_type == 1) & (distcarriercdrdf.busi_type == 1)]
    valuetuple = distcarrierin.percent.values
    valuetuplefirst5 = distcarrierin[distcarrierin.carrier_cd != 'Others'].percent.values

    # first5 percent
    incdrfirstcarrierpls = slide.placeholders[10]
    incdrfirstcarrierpls.text = str(valuetuplefirst5.sum()*100)+'%'
    # incdrfirstcarrierpls.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # valuetuple.resize(6)
    # otherpct = 1 - valuetuple.sum()
    # valuetuple[5]=otherpct
    # valuetuple = valuetuple.tolist().append(otherpct)
    # print(valuetuple)
    # print(list(distcarrierin.carrier_cd.values).append('Others'))
    # define chart data ---------------------
    chart_data = ChartData()
    # carrierCategories = ('HKGPP', 'TWNFE', 'MACCT', 'USACG', 'THAWP', 'Others')
    carrierCategories = distcarrierin.carrier_cd.values
    chart_data.categories = carrierCategories
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per prov
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distprovcdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_prov.xlsx")
    distprovin = distprovcdrdf[
        ['prov_cd', 'percent']][
        (str(distprovcdrdf.month.values[0]) == MONTH_ANA) & (distprovcdrdf.roam_type == 1) & (
            distprovcdrdf.busi_type == 1)]
    valuetuple = distprovin.percent.values
    valuetuplefirst5 = distprovin[distprovin.prov_cd != u'其他'].percent.values

    # first5 percent
    incdrfirstprovpls = slide.placeholders[11]
    incdrfirstprovpls.text = str(valuetuplefirst5.sum()*100)+'%'

    # define chart data
    chart_data = ChartData()
    provCategories = ('BJ', 'GD', 'SH', 'JS', 'ZJ', 'Others')
    # provCategories = distprovin.prov_cd.values
    chart_data.categories = provCategories

    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(chart_data, x, y, cx, cy)

    # slide1.1 来访话单量环比情况
    title_only_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    # slide1.1 chart1 per day
    title = ''
    titlenm = u"来访话单量按日情况"
    unit = u"单位：百万条"
    titlepls = slide.placeholders[13]
    titlepls.text = titlenm
    unitpls = slide.placeholders[12]
    unitpls.text = unit

    # valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0, 16.7, 19.2, 21.4, 16.7, 12.0)
    doddf = pd.read_excel("E:\djangoWeb\pptgen\dod_pptgen.xlsx")
    dodin = doddf[
        ['count', 'user', 'callduration', 'sms', 'dataall']][
        (str(doddf.date.values[0])[0:6] == MONTH_ANA) & (doddf.roam_type == 1)]
    valuetuple = dodin['count'].values
    # define chart data
    chart_data = ChartData()
    chart_data.categories = SILDE1x_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.7), Inches(10), Inches(2.8)
    createLineChart(chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per carrier
    # piechar
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distcarriercdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_carrier.xlsx")
    distcarrierin = distcarriercdrdf[
        ['carrier_cd', 'percent']][
        (str(distcarriercdrdf.month.values[0]) == MONTH_ANA) & (distcarriercdrdf.roam_type == 1) & (
        distcarriercdrdf.busi_type == 1)]
    valuetuple = distcarrierin.percent.values
    valuetuplefirst5 = distcarrierin[distcarrierin.carrier_cd != 'Others'].percent.values

    # first5 percent
    incdrfirstcarrierpls = slide.placeholders[10]
    incdrfirstcarrierpls.text = str(valuetuplefirst5.sum() * 100) + '%'
    # incdrfirstcarrierpls.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # valuetuple.resize(6)
    # otherpct = 1 - valuetuple.sum()
    # valuetuple[5]=otherpct
    # valuetuple = valuetuple.tolist().append(otherpct)
    # print(valuetuple)
    # print(list(distcarrierin.carrier_cd.values).append('Others'))
    # define chart data ---------------------
    chart_data = ChartData()
    # carrierCategories = ('HKGPP', 'TWNFE', 'MACCT', 'USACG', 'THAWP', 'Others')
    carrierCategories = distcarrierin.carrier_cd.values
    chart_data.categories = carrierCategories
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(chart_data, x, y, cx, cy)

    # slide1.1 chart1 pie per prov
    title = ""
    # valuetuple = (0.135, 0.324, 0.180, 0.135, 0.126, 0.1)
    distprovcdrdf = pd.read_excel("E:\djangoWeb\pptgen\\first5_pptgen_prov.xlsx")
    distprovin = distprovcdrdf[
        ['prov_cd', 'percent']][
        (str(distprovcdrdf.month.values[0]) == MONTH_ANA) & (distprovcdrdf.roam_type == 1) & (
            distprovcdrdf.busi_type == 1)]
    valuetuple = distprovin.percent.values
    valuetuplefirst5 = distprovin[distprovin.prov_cd != u'其他'].percent.values

    # first5 percent
    incdrfirstprovpls = slide.placeholders[11]
    incdrfirstprovpls.text = str(valuetuplefirst5.sum() * 100) + '%'

    # define chart data
    chart_data = ChartData()
    provCategories = ('BJ', 'GD', 'SH', 'JS', 'ZJ', 'Others')
    # provCategories = distprovin.prov_cd.values
    chart_data.categories = provCategories

    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(3.8), Inches(4.7), Inches(3.5)
    createPieChart(chart_data, x, y, cx, cy)

    # slide2 传统业务量趋势情况 - 来访
    title_only_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    # slide2 chart1 CDR
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.5), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide2 chart2 CallDuration
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(2.7), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide2 chart3 SMS
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(4.9), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide3 传统业务量趋势情况 - 出访
    title_only_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    # slide3 chart1 CDR
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(0.5), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide3 chart2 CallDuration
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(2.7), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide3 chart3 SMS
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series(title, valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0), Inches(4.9), Inches(10), Inches(2.5)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 数据业务量趋势情况
    title_only_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    # slide4 chart1 来访
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data1 ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("23g", valuetuple)
    chart_data.add_series("4g", valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0.1), Inches(1.2), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 chart2 出访
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data1 ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("23g", valuetuple)
    chart_data.add_series("4g", valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(1.2), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 chart3 用户数
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("In", valuetuple)
    chart_data.add_series("Out", valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(0.1), Inches(4.7), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 chart4 DOU
    title = ""
    valuetuple = (19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 19.2, 21.4, 16.7, 12.0)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = SILDE1_CATEGORIES
    chart_data.add_series("In", valuetuple)
    chart_data.add_series("Out", valuetuple)
    # 整体布局偏移量
    x, y, cx, cy = Inches(5), Inches(4.7), Inches(5), Inches(2.3)
    createTrendChart(chart_data, x, y, cx, cy)

    # slide4 paragraph
    inroam4gcarriersnuminc = 4
    inroam4gcarriersnumall = 103
    outroam4gcarriersnuminc = 4
    outroam4gcarriersnumall = 122
    para1 = slide.placeholders[15]
    para1.text = "2016年02月来访LTE新开通" + str(inroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(inroam4gcarriersnumall) + "家；" + "\r" + "2016年02月出访LTE新开通" + str(outroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(outroam4gcarriersnumall) + "家；"
    para2 = slide.placeholders[16]
    para2.text = "来访：4%   出访：-3%"
    para3 = slide.placeholders[17]
    para3.text = "来访：4%   出访：-3%"

    # tf = para.text_frame
    # p = tf.add_paragraph()
    # p.text = "2016年02月来访LTE新开通" + str(inroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(inroam4gcarriersnumall) + "家；"
    # p.level = 1
    #
    # p = tf.add_paragraph()
    # p.text = "2016年02月出访LTE新开通" + str(outroam4gcarriersnuminc) + "家运营商（MYSMT、 MACSM等），累计" + str(outroam4gcarriersnumall) + "家；"
    # p.level = 1

    prs.save('report1.pptx')
