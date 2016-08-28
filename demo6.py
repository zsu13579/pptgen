# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(12.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(2.0)
table.columns[4].width = Inches(2.0)
table.columns[5].width = Inches(2.0)

# write column headings
table.cell(0, 0).text = '漫游方向'
table.cell(0, 1).text = '总话单量（亿条）'
table.cell(0, 2).text = '总用户数（百万户）'
table.cell(0, 3).text = '通话时长（千万分钟）'
table.cell(0, 4).text = '短信条数（亿条）'
table.cell(0, 5).text = '流量（TB）'

# write body cells
table.cell(1, 0).text = '来访'
table.cell(1, 1).text = '1'
table.cell(1, 2).text = '2'
table.cell(1, 3).text = '3.5'
table.cell(1, 4).text = '4'
table.cell(1, 5).text = '98'

# write body cells
table.cell(2, 0).text = '出访'
table.cell(2, 1).text = '3'
table.cell(2, 2).text = '1'
table.cell(2, 3).text = '2.5'
table.cell(2, 4).text = '3'
table.cell(2, 5).text = '92'

prs.save('report1.pptx')