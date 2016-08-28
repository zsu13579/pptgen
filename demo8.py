#-*- coding:utf-8 -*-
from pptx import Presentation

def analyze_ppt(input, output):
    """ 读入文件并分析结构。
    输出文件包含标记信息，使生成后面的PowerPoint模板更容易。
    """
    prs = Presentation(input)
    # 每个PowerPoint文件有多种布局
    # 循环找出不同的元素位置
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # 不是每张幻灯片都有标题
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # 遍历所有占位符，并通过索引和类型识别它们
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # 不要覆写仅仅是一个特殊占位符的标题
                try:
                    if 'Title' not in shape.text:
                        shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)

analyze_ppt("ana-template.pptx", "ana-template-markup.pptx")
